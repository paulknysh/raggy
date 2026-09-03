import logging

import pytest
from langchain_core.documents import Document

from raggy import loaders
from raggy.loaders import load_documents


def test_load_documents_loads_multiple_sources(tmp_path):
    dir_a = tmp_path / "a"
    dir_a.mkdir()
    (dir_a / "one.txt").write_text("Text A", encoding="utf-8")
    (tmp_path / "two.md").write_text("Text B", encoding="utf-8")

    documents = load_documents([str(dir_a), str(tmp_path / "two.md")])

    contents = {doc.page_content for doc in documents}
    assert contents == {"Text A", "Text B"}


def test_load_documents_raises_for_missing_source(tmp_path):
    dir_a = tmp_path / "a"
    dir_a.mkdir()
    (dir_a / "one.txt").write_text("Text A", encoding="utf-8")

    with pytest.raises(FileNotFoundError, match="not found"):
        load_documents([str(dir_a), str(tmp_path / "missing.txt")])


def test_load_documents_ignores_unsupported_extension(tmp_path):
    unsupported = tmp_path / "notes.csv"
    unsupported.write_text("hello", encoding="utf-8")

    documents = load_documents([str(unsupported)])

    assert documents == []


def test_load_documents_reports_unsupported_file_once(tmp_path, caplog):
    (tmp_path / "a.txt").write_text("Text A", encoding="utf-8")
    (tmp_path / "notes.csv").write_text("ignored", encoding="utf-8")
    (tmp_path / "data.json").write_text("{}", encoding="utf-8")

    with caplog.at_level(logging.INFO, logger="raggy.loaders"):
        documents = load_documents([str(tmp_path)])

    contents = {doc.page_content for doc in documents}
    assert contents == {"Text A"}

    unsupported_messages = [
        record.getMessage()
        for record in caplog.records
        if "unsupported" in record.getMessage().lower()
    ]
    assert len(unsupported_messages) == 1
    assert "2" in unsupported_messages[0]


def test_load_documents_loads_markdown_file(tmp_path):
    doc = tmp_path / "notes.md"
    doc.write_text("# Heading\n\nSome markdown body.", encoding="utf-8")

    documents = load_documents([str(doc)])

    assert len(documents) == 1
    assert "Some markdown body." in documents[0].page_content


def test_load_documents_loads_text_file(tmp_path):
    doc = tmp_path / "notes.txt"
    doc.write_text("Hello world", encoding="utf-8")

    documents = load_documents([str(doc)])

    assert len(documents) == 1
    assert documents[0].page_content == "Hello world"


def test_load_documents_loads_pdf_file(tmp_path):
    import pymupdf

    pdf_path = tmp_path / "sample.pdf"
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 100), "PDF content here")
    doc.save(str(pdf_path))

    documents = load_documents([str(pdf_path)])

    assert len(documents) == 1
    assert "PDF content here" in documents[0].page_content
    assert documents[0].metadata["page"] == 1


def test_load_documents_pdf_pages_are_one_indexed(tmp_path):
    import pymupdf

    pdf_path = tmp_path / "pages.pdf"
    doc = pymupdf.open()
    for _ in range(3):
        page = doc.new_page()
        page.insert_text((72, 100), "Page content")
    doc.save(str(pdf_path))

    documents = load_documents([str(pdf_path)])

    assert [doc.metadata["page"] for doc in documents] == [1, 2, 3]


def test_load_documents_recursively_loads_directory(tmp_path):
    (tmp_path / "nested").mkdir()
    (tmp_path / "a.txt").write_text("Text A", encoding="utf-8")
    (tmp_path / "nested" / "b.txt").write_text("Text B", encoding="utf-8")
    (tmp_path / "ignore.csv").write_text("ignored", encoding="utf-8")

    documents = load_documents([str(tmp_path)])

    contents = {doc.page_content for doc in documents}
    assert contents == {"Text A", "Text B"}


def test_load_documents_ignores_directory_with_no_supported_files(tmp_path, caplog):
    """An empty source contributes nothing; only a missing one is an error."""
    (tmp_path / "notes.csv").write_text("ignored", encoding="utf-8")

    with caplog.at_level(logging.WARNING, logger="raggy.loaders"):
        documents = load_documents([str(tmp_path)])

    assert documents == []
    # Nothing loaded at all, so the "unsupported" notice is raised to WARNING,
    # which is the level the CLI actually shows.
    assert any(
        record.levelno == logging.WARNING and "unsupported" in record.getMessage()
        for record in caplog.records
    )


def test_supported_extensions_exported():
    assert loaders.SUPPORTED_EXTENSIONS == {
        ".txt",
        ".md",
        ".markdown",
        ".pdf",
        ".docx",
        ".pptx",
        ".html",
        ".htm",
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
    }


def test_load_documents_loads_docx_file(tmp_path):
    from docx import Document as DocxDocument

    docx_path = tmp_path / "notes.docx"
    docx = DocxDocument()
    docx.add_paragraph("Hello from docx")
    docx.save(str(docx_path))

    documents = load_documents([str(docx_path)])

    assert len(documents) == 1
    assert "Hello from docx" in documents[0].page_content
    assert "page" not in documents[0].metadata


def test_load_documents_loads_pptx_file(tmp_path):
    from pptx import Presentation

    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide Title"
    presentation.save(str(pptx_path))

    documents = load_documents([str(pptx_path)])

    assert len(documents) == 1
    assert "Slide Title" in documents[0].page_content
    assert documents[0].metadata["page"] == 1


def test_load_documents_loads_html_file(tmp_path):
    html_path = tmp_path / "page.html"
    html_path.write_text(
        "<html><body><h1>Heading</h1><p>Some HTML body.</p></body></html>",
        encoding="utf-8",
    )

    documents = load_documents([str(html_path)])

    assert len(documents) == 1
    assert "Some HTML body." in documents[0].page_content


def test_load_documents_directory_includes_new_formats(tmp_path):
    from docx import Document as DocxDocument

    docx = DocxDocument()
    docx.add_paragraph("Docx content")
    docx.save(str(tmp_path / "report.docx"))
    (tmp_path / "page.html").write_text(
        "<html><body>HTML content</body></html>", encoding="utf-8"
    )

    documents = load_documents([str(tmp_path)])

    contents = " ".join(doc.page_content for doc in documents)
    assert "Docx content" in contents
    assert "HTML content" in contents


def _make_ocr_image(path) -> None:
    from PIL import Image, ImageDraw, ImageFont

    font = None
    for candidate in (
        "/System/Library/Fonts/Helvetica.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ):
        try:
            font = ImageFont.truetype(candidate, 48)
            break
        except OSError:
            continue
    if font is None:
        font = ImageFont.load_default()
    img = Image.new("RGB", (1000, 300), "white")
    draw = ImageDraw.Draw(img)
    draw.text((20, 30), "HELLO WORLD", fill="black", font=font)
    draw.text((20, 110), "RAGGY PIPELINE", fill="black", font=font)
    img.save(str(path))


def _compact(text: str) -> str:
    return "".join(text.split())


def test_load_documents_loads_image_file(tmp_path):
    image_path = tmp_path / "scan.png"
    _make_ocr_image(image_path)

    documents = load_documents([str(image_path)])

    assert len(documents) == 1
    assert "HELLOWORLD" in _compact(documents[0].page_content)


def test_load_documents_ocr_fallback_for_scanned_pdf(tmp_path):
    import pymupdf

    image_path = tmp_path / "scan.png"
    _make_ocr_image(image_path)

    pdf_path = tmp_path / "scanned.pdf"
    doc = pymupdf.open()
    page = doc.new_page(width=1000, height=300)
    page.insert_image(pymupdf.Rect(0, 0, 1000, 300), filename=str(image_path))
    doc.save(str(pdf_path))

    documents = load_documents([str(pdf_path)])

    assert len(documents) == 1
    assert "HELLOWORLD" in _compact(documents[0].page_content)
    assert documents[0].metadata["page"] == 1


def test_load_documents_loads_only_the_given_files(tmp_path):
    first = tmp_path / "first.txt"
    second = tmp_path / "second.txt"
    ignored = tmp_path / "ignored.txt"
    first.write_text("first content", encoding="utf-8")
    second.write_text("second content", encoding="utf-8")
    ignored.write_text("ignored content", encoding="utf-8")

    documents = loaders.load_documents([first, second])

    assert [doc.page_content for doc in documents] == [
        "first content",
        "second content",
    ]
    assert [doc.metadata["source"] for doc in documents] == [str(first), str(second)]


def test_load_documents_skips_missing_when_on_missing_is_skip(tmp_path):
    present = tmp_path / "present.md"
    present.write_text("kept", encoding="utf-8")
    unsupported = tmp_path / "data.csv"
    unsupported.write_text("a,b", encoding="utf-8")

    documents = loaders.load_documents(
        [present, unsupported, tmp_path / "gone.txt"], on_missing="skip"
    )

    assert [doc.page_content for doc in documents] == ["kept"]


def test_load_documents_raises_for_missing_file_by_default(tmp_path):
    """The incremental caller opts into skipping; nobody else gets it silently."""
    present = tmp_path / "present.md"
    present.write_text("kept", encoding="utf-8")

    with pytest.raises(FileNotFoundError, match="not found"):
        loaders.load_documents([present, tmp_path / "gone.txt"])


def test_load_documents_rejects_an_unknown_on_missing(tmp_path):
    with pytest.raises(ValueError, match="on_missing"):
        loaders.load_documents([tmp_path], on_missing="ignore")


def test_load_documents_reports_progress_per_file(tmp_path):
    dir_a = tmp_path / "a"
    dir_a.mkdir()
    (dir_a / "one.txt").write_text("Text A", encoding="utf-8")
    (dir_a / "skip.csv").write_text("ignored", encoding="utf-8")
    (tmp_path / "two.md").write_text("Text B", encoding="utf-8")

    messages: list[str] = []
    load_documents([str(dir_a), str(tmp_path / "two.md")], progress=messages.append)

    assert messages == [
        "[1/2] ingesting one.txt ...",
        "[2/2] ingesting two.md ...",
    ]


def test_load_documents_progress_excludes_missing_files(tmp_path):
    first = tmp_path / "first.txt"
    first.write_text("first", encoding="utf-8")
    missing = tmp_path / "gone.txt"

    messages: list[str] = []
    loaders.load_documents(
        [first, missing], progress=messages.append, on_missing="skip"
    )

    assert messages == ["[1/1] ingesting first.txt ..."]


def test_loading_walks_the_sources_exactly_once(tmp_path, monkeypatch):
    """One walk feeds both the loaded files and the i/N denominator."""
    (tmp_path / "one.txt").write_text("Text A", encoding="utf-8")

    walks = []
    real_walk = loaders._walk_source
    monkeypatch.setattr(
        loaders,
        "_walk_source",
        lambda root, skipped: walks.append(root) or real_walk(root, skipped),
    )

    assert load_documents([str(tmp_path)], progress=lambda *a, **k: None)
    assert walks == [tmp_path]


def test_overlapping_sources_load_each_file_once(tmp_path):
    """The walk backing the loaders and the manifest must agree on the file set.

    A file reachable through two source entries has one fingerprint, so it must
    also produce one set of chunks — otherwise incremental indexing tracks a
    corpus the DB does not contain.
    """
    from raggy.vectorstore import file_fingerprints

    nested = tmp_path / "sub"
    nested.mkdir()
    (tmp_path / "a.txt").write_text("Text A", encoding="utf-8")
    (nested / "b.txt").write_text("Text B", encoding="utf-8")

    sources = [str(tmp_path), str(nested), str(nested / "b.txt")]

    files = loaders.source_files(sources)
    documents = load_documents(sources)
    fingerprints = file_fingerprints(sources)

    assert [str(path) for path in files] == sorted(fingerprints)
    assert [doc.metadata["source"] for doc in documents] == [
        str(path) for path in files
    ]
    assert [doc.page_content for doc in documents] == ["Text A", "Text B"]


def test_annotate_line_numbers_sets_start_and_end_lines():
    content = "line one\nline two\nline three\nline four\nline five\n"

    split = Document(
        page_content="line three\nline four\n",
        metadata={"source": "doc.txt"},
    )

    loaders.annotate_line_numbers([split], content)

    assert split.metadata["start_line"] == 3
    assert split.metadata["end_line"] == 4


def test_annotate_line_numbers_counts_first_line_as_one():
    content = "single line with no trailing newline"

    split = Document(page_content="single line with", metadata={})

    loaders.annotate_line_numbers([split], content)

    assert split.metadata["start_line"] == 1
    assert split.metadata["end_line"] == 1


def test_annotate_line_numbers_falls_back_for_overlapping_chunks():
    content = "alpha\nbeta\ngamma\nalpha\nbeta\n"

    split = Document(page_content="alpha\nbeta\n", metadata={})

    loaders.annotate_line_numbers([split], content)

    assert split.metadata["start_line"] == 1
    assert split.metadata["end_line"] == 2


def test_should_annotate_lines_only_for_verbatim_text_files():
    line_annotated = [".txt", ".md", ".markdown"]
    for suffix in line_annotated:
        doc = Document(page_content="x", metadata={"source": f"/tmp/doc{suffix}"})
        assert loaders.should_annotate_lines(doc) is True, suffix

    # .html/.htm are excluded because BSHTMLLoader returns extracted text, not
    # the file, so a line counted in it need not be that line of the source.
    locationless = [
        ".html",
        ".htm",
        ".pdf",
        ".pptx",
        ".docx",
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
    ]
    for suffix in locationless:
        doc = Document(page_content="x", metadata={"source": f"/tmp/doc{suffix}"})
        assert loaders.should_annotate_lines(doc) is False, suffix
