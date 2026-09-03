import logging
from collections.abc import Iterator, Sequence
from pathlib import Path
from typing import Literal

from langchain_community.document_loaders import (
    BSHTMLLoader,
    Docx2txtLoader,
    PyPDFLoader,
    TextLoader,
)
from langchain_core.documents import Document

from .progress import ProgressCallback

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".pdf",
    ".docx",
    ".pptx",
    ".html",
    ".htm",
    ".png",
    ".jpg",
    ".jpeg",
    ".bmp",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp"}

# The file types whose chunks get start_line/end_line annotations. Only
# formats whose loader returns the file verbatim qualify; see
# should_annotate_lines for why the rest carry no location metadata.
LINE_ANNOTATED_EXTENSIONS = {".txt", ".md", ".markdown"}

DEFAULT_OCR_DPI = 150


_ocr_engine = None


def _get_ocr_engine():
    """Return a lazily-initialized shared RapidOCR engine instance."""
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR

        _ocr_engine = RapidOCR()
    return _ocr_engine


def _ocr_image_bytes(image_bytes: bytes) -> str:
    """Run OCR on raw image bytes and return the recognized text."""
    result, _ = _get_ocr_engine()(image_bytes)
    if not result:
        return ""
    return "\n".join(line[1] for line in result if line[1])


def _load_image(path: Path) -> list[Document]:
    """OCR a single image file into one Document."""
    text = _ocr_image_bytes(path.read_bytes())
    return [Document(page_content=text, metadata={"source": str(path)})]


def _load_ocr_pdf(path: Path) -> list[Document]:
    """OCR an image-only PDF by rendering each page and running OCR per page."""
    import pymupdf

    documents: list[Document] = []
    with pymupdf.open(str(path)) as pdf:
        for page_number in range(pdf.page_count):
            page = pdf[page_number]
            pixmap = page.get_pixmap(dpi=DEFAULT_OCR_DPI)
            text = _ocr_image_bytes(pixmap.tobytes("png"))
            documents.append(
                Document(
                    page_content=text,
                    metadata={"source": str(path), "page": page_number + 1},
                )
            )
    return documents


def _collect_text_frame(frame, text_parts: list[str]) -> None:
    """Append the text of every non-empty paragraph in a text frame."""
    text = "\n".join(
        paragraph.text for paragraph in frame.paragraphs if paragraph.text.strip()
    )
    if text.strip():
        text_parts.append(text)


def _collect_shape(shape, text_parts: list[str]) -> None:
    """Recursively collect text from a shape, its text frames, and its tables."""
    if getattr(shape, "has_text_frame", False):
        _collect_text_frame(shape.text_frame, text_parts)
    elif getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                text_parts.append(" | ".join(cells))
    elif getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            _collect_shape(child, text_parts)


def _load_pptx(path: Path) -> list[Document]:
    """Extract text from every slide of a PowerPoint deck, one Document per slide."""
    from pptx import Presentation

    presentation = Presentation(str(path))
    documents: list[Document] = []

    for index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []

        for shape in slide.shapes:
            _collect_shape(shape, text_parts)

        if text_parts:
            documents.append(
                Document(
                    page_content="\n\n".join(text_parts),
                    metadata={"source": str(path), "page": index},
                )
            )

    return documents


def _load_file(path: Path) -> list[Document]:
    """Load documents from a single supported file based on its extension.

    Callers reach this only through :func:`source_files`, which is what decides
    a file is supported, so the extension is dispatched on here but never
    re-checked.
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        loader = PyPDFLoader(str(path))
        documents = loader.load()
        for doc in documents:
            if "page" in doc.metadata:
                doc.metadata["page"] += 1
        if not any((doc.page_content or "").strip() for doc in documents):
            logger.info("No extractable text in '%s'; falling back to OCR...", path)
            documents = _load_ocr_pdf(path)
        return documents

    if suffix in IMAGE_EXTENSIONS:
        return _load_image(path)

    if suffix == ".docx":
        return Docx2txtLoader(str(path)).load()

    if suffix == ".pptx":
        return _load_pptx(path)

    if suffix in {".html", ".htm"}:
        return BSHTMLLoader(str(path), open_encoding="utf-8").load()

    loader = TextLoader(str(path), encoding="utf-8")
    return loader.load()


def annotate_line_numbers(splits: list[Document], content: str) -> None:
    """Annotate text splits with ``start_line``/``end_line`` line ranges.

    The positions are derived from where each chunk's text appears in the
    ``content`` string (1-indexed line numbers). Because ``chunk_overlap``
    causes adjacent chunks to share text, a search from ``cursor`` may not
    find a chunk against its true start; the first-occurrence fallback still
    yields approximate (but useful) line attribution.
    """
    cursor = 0
    chunk_size_tolerance = max(0, len(content))
    for split in splits:
        text = split.page_content
        start = content.find(text, cursor)
        if start == -1:
            start = content.find(text, 0, cursor + chunk_size_tolerance)
        if start == -1:
            start = content.find(text)

        if start == -1:
            continue
        split.metadata["start_line"] = content.count("\n", 0, start) + 1
        if text.endswith("\n"):
            newlines_in_text = text.count("\n") - 1
        else:
            newlines_in_text = text.count("\n")
        split.metadata["end_line"] = split.metadata["start_line"] + newlines_in_text
        cursor = start + len(text)


def should_annotate_lines(doc: Document) -> bool:
    """Return True only for text-based files that get line-number annotations.

    Line numbers are counted in the loaded ``page_content``, so they describe
    the file itself only for formats whose loader returns it verbatim — the
    plain-text ones, read by ``TextLoader``. HTML is excluded for exactly this
    reason: ``BSHTMLLoader`` yields the extracted text, which keeps the
    newlines inside text nodes but drops those inside tags and comments, so
    the count drifts further from the real line the deeper into the file a
    chunk sits.

    PDFs and PowerPoint decks carry a ``page`` key set by their loaders. DOCX
    has no native page boundaries (and Word's pagination can't be reproduced
    reliably), and image files have no meaningful lines, so both carry no
    location metadata and are skipped here.
    """
    suffix = Path(doc.metadata.get("source", "")).suffix.lower()
    return suffix in LINE_ANNOTATED_EXTENSIONS


def _walk_source(root: Path, skipped: list[Path]) -> Iterator[Path]:
    """Yield the supported files one source entry contributes, in load order."""
    if root.is_dir():
        for file_path in sorted(root.rglob("*")):
            if not file_path.is_file():
                continue
            if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
                skipped.append(file_path)
                continue
            yield file_path
    elif root.suffix.lower() in SUPPORTED_EXTENSIONS:
        yield root
    else:
        skipped.append(root)


def source_files(
    sources: Sequence[str | Path],
    skipped: list[Path] | None = None,
    on_missing: Literal["raise", "skip"] = "raise",
) -> list[Path]:
    """Return every supported file under ``sources``, in the order it is indexed.

    This is the corpus's single walk: the loaders below read the files it names
    and :func:`raggy.vectorstore.file_fingerprints` hashes them, so the paths
    recorded in the manifest are by construction the ones the loaders write
    into each chunk's ``source`` metadata.

    Each entry may be a single file or a directory, which is walked
    recursively. Unsupported files are recorded in ``skipped`` rather than
    raising, and a directory holding none simply contributes nothing. A file
    reachable through more than one entry is returned once, at its first
    position. ``on_missing`` decides what an entry that no longer exists means
    (see :func:`load_documents`).
    """
    if on_missing not in ("raise", "skip"):
        raise ValueError(f"on_missing must be 'raise' or 'skip', got {on_missing!r}")

    paths: list[Path] = []
    missing: list[Path] = []
    for source in sources:
        path = Path(source)
        (paths if path.exists() else missing).append(path)

    if missing and on_missing == "raise":
        raise FileNotFoundError(
            "Source document(s) not found at: "
            + ", ".join(str(path) for path in missing)
        )
    for path in missing:
        logger.warning("Skipping '%s': file no longer exists.", path)

    files: list[Path] = []
    seen: set[Path] = set()

    def claim(path: Path) -> bool:
        """True the first time a path is reached, under any of its spellings."""
        resolved = path.resolve()
        if resolved in seen:
            return False
        seen.add(resolved)
        return True

    for root in paths:
        unsupported: list[Path] = []
        for file_path in _walk_source(root, unsupported):
            if claim(file_path):
                files.append(file_path)
        if skipped is not None:
            skipped.extend(path for path in unsupported if claim(path))
    return files


def _load_each(
    paths: list[Path],
    progress: ProgressCallback | None,
) -> list[Document]:
    """Load every path in turn, reporting each one and skipping what fails.

    A file that cannot be read is logged and passed over rather than aborting
    the run: one unreadable file should not cost an otherwise good corpus.
    """
    total = len(paths)
    documents: list[Document] = []
    for index, path in enumerate(paths, start=1):
        if progress is not None:
            progress(f"[{index}/{total}] ingesting {path.name} ...")
        try:
            documents.extend(_load_file(path))
        except Exception as e:  # noqa: BLE001
            logger.warning("Failed to load '%s': %s", path, e)
    return documents


def _report_unsupported(skipped: list[Path], loaded_any: bool) -> None:
    """Log a single message summarizing ignored unsupported files.

    Raised to WARNING when nothing loaded at all: the sources then hold nothing
    readable — most likely a directory of file types raggy does not support —
    and the CLI shows WARNING but not INFO.
    """
    if not skipped:
        return
    logger.log(
        logging.INFO if loaded_any else logging.WARNING,
        "Detected %d unsupported file(s) that won't be used; ignoring them. "
        "Supported file types: %s.",
        len(skipped),
        ", ".join(sorted(SUPPORTED_EXTENSIONS)),
    )


def load_documents(
    sources: Sequence[str | Path],
    progress: ProgressCallback | None = None,
    on_missing: Literal["raise", "skip"] = "raise",
) -> list[Document]:
    """Load documents from files and/or directories.

    Each entry may be a single supported file or a directory containing them;
    see :func:`source_files` for how the two are walked. Unsupported file types
    are ignored and reported once, and a file that fails to read is logged and
    skipped rather than aborting the run. ``progress`` receives one status line
    per file, counted across all entries.

    ``on_missing`` decides what an entry that no longer exists means. The
    default ``"raise"`` suits the configured sources, where a missing path is a
    config error worth stopping for. Incremental indexing passes ``"skip"``:
    its file list was fingerprinted earlier in the same run, so a file deleted
    since then is expected, and dropping it beats aborting the whole update —
    the next run reconciles it as a deletion.
    """
    skipped: list[Path] = []
    files = source_files(sources, skipped, on_missing=on_missing)
    documents = _load_each(files, progress)
    _report_unsupported(skipped, bool(documents))
    return documents
